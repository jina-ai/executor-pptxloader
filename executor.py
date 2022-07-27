from typing import Optional, Dict
import pptx
from warnings import warn
from jina import Executor, DocumentArray, requests, Document


class PptxLoader(Executor):
    """An Executor for loading text and images from Powerpoint .pptx files"""
    def __init__(self,
                 access_paths: str = 'r',
                 *args, **kwargs):
        """
        :param access_paths: the access paths to be used in traversing
            the DocumentArray received
        :param args: the *args for Executor
        :param kwargs: the **kwargs for Executor
        """
        if("traversal_paths" in kwargs.keys()):
            warn("'traversal_paths' is deprecated, please use 'access_paths'")
        super().__init__(*args, **kwargs)
        self.access_paths = access_paths

    @requests
    def process(self, docs: DocumentArray, parameters: Optional[Dict], **kwargs):
        """Process the documents and extract the text and images
        """
        access_paths = parameters.get('access_paths', self.access_paths)
        for d in docs.traverse_flat(access_paths):
            self._process_one(d)

    def _process_one(self, d):
        prs = pptx.Presentation(d.uri)
        nr = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            d.chunks.append(Document(content=run.text,
                                                     modality='text'))

                if isinstance(shape, pptx.shapes.picture.Picture):
                    d.chunks.append(Document(content=shape.image.blob,
                                             modality='image'))
                    nr += 1

